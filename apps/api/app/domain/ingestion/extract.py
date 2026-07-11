"""Extração de texto por família de formato (ADR-0002).

Quando `settings.vision_enabled` (ADR-0012), imagens embutidas em PDF/DOCX são
descritas por LLM vision e intercaladas no texto antes do chunking — de forma
transparente para o pipeline (a assinatura de `extract_text` não muda).
"""
from __future__ import annotations

import logging
import os

from app.config import settings
from app.db.models import IMAGE_EXTENSIONS
from app.domain.ingestion.errors import PermanentIngestionError

log = logging.getLogger("worker.extract")

_TEXT_EXT = {".txt", ".md", ".py"}
_IMAGE_EXT = set(IMAGE_EXTENSIONS)


def extract_text(path: str, filename: str) -> str:
    fname = filename or os.path.basename(path)
    ext = os.path.splitext(fname)[1].lower()
    try:
        if ext in _TEXT_EXT:
            return _read_text(path)
        if ext in {".html", ".htm"}:
            return _extract_html(path)
        if ext == ".pdf":
            return _extract_pdf(path, fname)
        if ext == ".docx":
            return _extract_docx(path, fname)
        if ext in {".xls", ".xlsx"}:
            return _extract_xlsx(path)
        if ext == ".pptx":
            return _extract_pptx(path, fname)
        if ext == ".ppt":  # formato legado (OLE) não lido pelo python-pptx
            raise PermanentIngestionError("Formato .ppt legado não suportado — converta para .pptx.")
        if ext == ".ipynb":
            return _extract_ipynb(path)
        if ext in _IMAGE_EXT:
            return _extract_image(path, fname)
    except PermanentIngestionError:
        raise
    except Exception as exc:  # parser falhou → arquivo provavelmente inválido (permanente)
        raise PermanentIngestionError(f"Falha ao extrair {ext}: {exc}") from exc
    raise PermanentIngestionError(f"Extensão não suportada: {ext}")


def _read_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _extract_html(path: str) -> str:
    from bs4 import BeautifulSoup

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return soup.get_text(separator="\n")


def _extract_pdf(path: str, filename: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)

    def _text_only() -> str:
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if not settings.vision_enabled:
        return _text_only()
    try:
        return _extract_pdf_with_vision(path, filename, reader)
    except Exception as exc:  # noqa: BLE001 — vision é best-effort (ADR-0012): degrada p/ texto
        log.warning("Extração com vision falhou em %s (segue só com texto): %s", filename, exc)
        return _text_only()


def _extract_pdf_with_vision(path: str, filename: str, reader) -> str:
    """Texto via pypdf; imagens embutidas via PyMuPDF (ADR-0012), intercaladas por página.

    Best-effort por página: falha ao abrir/ler uma página não interrompe as demais.
    Deduplica por `xref` — imagem repetida (ex.: logo) reusa a descrição, sem nova chamada.
    """
    import fitz  # PyMuPDF

    from app.services import vision

    parts: list[str] = []
    desc_by_xref: dict[int, str] = {}
    fdoc = fitz.open(path)
    try:
        for page_num, page in enumerate(reader.pages, start=1):
            parts.append(page.extract_text() or "")
            if page_num - 1 >= fdoc.page_count:
                continue
            try:
                fpage = fdoc.load_page(page_num - 1)
                images = fpage.get_images(full=True)
            except Exception:  # noqa: BLE001 — página problemática não interrompe as demais
                continue
            for img in images:
                xref = img[0]
                if xref not in desc_by_xref:
                    image_bytes = _image_png_for_vision(fdoc, fpage, xref)
                    desc_by_xref[xref] = (
                        vision.describe_image(image_bytes, filename, f"página {page_num}")
                        if image_bytes
                        else ""
                    )
                desc = desc_by_xref[xref]
                if desc:
                    parts.append(f"[Imagem — página {page_num}: {desc}]")
    finally:
        fdoc.close()
    return "\n".join(parts)


def _image_png_for_vision(fdoc, fpage, xref) -> bytes | None:
    """PNG da imagem para o vision (ADR-0012). Sempre retorna PNG (casa com o MIME enviado).

    Renderiza a região onde a imagem aparece na página em alta DPI
    (`settings.vision_render_dpi`): normaliza colorspace, captura a imagem como
    exibida e garante tamanho legível — melhor para tabelas/prints densos que o
    bitmap embutido. Fallback: bitmap embutido convertido para PNG (RGB) se não
    houver retângulo de posição.
    """
    import fitz  # PyMuPDF

    try:
        rects = fpage.get_image_rects(xref)
        if rects:
            pix = fpage.get_pixmap(clip=rects[0], dpi=settings.vision_render_dpi)
        else:
            pix = fitz.Pixmap(fdoc, xref)
            if pix.n - pix.alpha >= 4:  # CMYK/separação → RGB para gerar PNG válido
                pix = fitz.Pixmap(fitz.csRGB, pix)
        return pix.tobytes("png")
    except Exception:  # noqa: BLE001 — imagem problemática não interrompe a extração
        return None


def _extract_docx(path: str, filename: str) -> str:
    import docx

    doc = docx.Document(path)
    if not settings.vision_enabled:
        return "\n".join(p.text for p in doc.paragraphs)
    try:
        return _extract_docx_with_vision(doc, filename)
    except Exception as exc:  # noqa: BLE001 — vision é best-effort (ADR-0012): degrada p/ texto
        log.warning("Extração com vision falhou em %s (segue só com texto): %s", filename, exc)
        return "\n".join(p.text for p in doc.paragraphs)


def _extract_docx_with_vision(doc, filename: str) -> str:
    """Texto via python-docx; imagens (blips do XML) descritas por vision (ADR-0012)."""
    from docx.oxml.ns import qn

    from app.services import vision

    part = doc.part
    parts: list[str] = []
    for ordinal, para in enumerate(doc.paragraphs, start=1):
        parts.append(para.text)
        for blip in para._p.findall(".//" + qn("a:blip")):
            rid = blip.get(qn("r:embed"))
            if not rid:
                continue
            try:
                image_bytes = part.related_parts[rid].blob
            except KeyError:  # relação ausente/externa — ignora
                continue
            desc = vision.describe_image(image_bytes, filename, f"parágrafo {ordinal}")
            if desc:
                parts.append(f"[Imagem — parágrafo {ordinal}: {desc}]")
    return "\n".join(parts)


def _extract_xlsx(path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def _extract_pptx(path: str, filename: str) -> str:
    """Texto por slide (título/corpo/tabelas + notas). Imagens dos slides via vision best-effort (ADR-0012)."""
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {idx}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
            if settings.vision_enabled:
                desc = _describe_pptx_picture(shape, filename, idx)
                if desc:
                    parts.append(f"[Imagem — slide {idx}: {desc}]")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notas do slide {idx}: {notes}]")
    return "\n".join(parts)


def _describe_pptx_picture(shape, filename: str, slide_num: int) -> str:
    """Descrição de uma imagem de slide — best-effort (só shapes do tipo Picture têm `.image`)."""
    from app.services import vision

    try:
        blob = shape.image.blob  # AttributeError em shapes que não são imagem
    except (AttributeError, Exception):  # noqa: BLE001 — imagem problemática não interrompe o slide
        return ""
    return vision.describe_image(blob, filename, f"slide {slide_num}")


def _extract_ipynb(path: str) -> str:
    """Notebook Jupyter: células markdown + código (com saídas textuais), na ordem do notebook."""
    import json

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        nb = json.load(f)

    def _joined(src) -> str:
        return ("".join(src) if isinstance(src, list) else (src or "")).strip()

    parts: list[str] = []
    for cell in nb.get("cells", []):
        ctype = cell.get("cell_type")
        src = _joined(cell.get("source", ""))
        if ctype == "markdown" and src:
            parts.append(src)
        elif ctype == "code":
            if src:
                parts.append(f"```\n{src}\n```")
            for out in cell.get("outputs", []):
                text = _joined(out.get("text", ""))  # stream
                if not text:
                    text = _joined(out.get("data", {}).get("text/plain", ""))  # execute_result/display_data
                if text:
                    parts.append(f"[Saída]: {text}")
    return "\n\n".join(parts)


def _extract_image(path: str, filename: str) -> str:
    """Imagem como documento próprio: descrição textual via vision (ADR-0012).

    Sem vision (ou em falha), degrada para um placeholder com o nome do arquivo — o
    documento ainda é indexável, apenas sem conteúdo semântico da imagem.
    """
    if not settings.vision_enabled:
        return f"[Imagem: {filename}]"
    from app.services import vision

    png = _image_file_to_png(path)
    desc = vision.describe_image(png, filename, "imagem") if png else ""
    return f"[Imagem — {filename}: {desc}]" if desc else f"[Imagem: {filename}]"


def _image_file_to_png(path: str) -> bytes | None:
    """Converte um arquivo de imagem para PNG (casa com o MIME enviado ao vision). None em falha."""
    import fitz  # PyMuPDF

    try:
        pix = fitz.Pixmap(path)
        if pix.n - pix.alpha >= 4:  # CMYK/separação → RGB para gerar PNG válido
            pix = fitz.Pixmap(fitz.csRGB, pix)
        return pix.tobytes("png")
    except Exception:  # noqa: BLE001 — imagem ilegível vira placeholder no chamador
        return None
