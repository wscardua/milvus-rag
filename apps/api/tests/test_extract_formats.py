"""WORK-007: extração dos novos formatos — .ipynb, imagem e PPTX (ADR-0002 rev.)."""
from __future__ import annotations

import json
import os
import tempfile

from app.config import settings
from app.domain.ingestion import extract


def _write(suffix: str, content: bytes) -> str:
    fd, path = tempfile.mkstemp(suffix=suffix)
    with os.fdopen(fd, "wb") as f:
        f.write(content)
    return path


def test_extract_ipynb_markdown_and_code():
    nb = {
        "cells": [
            {"cell_type": "markdown", "source": ["# Título\n", "texto md"]},
            {"cell_type": "code", "source": "print('oi')", "outputs": [
                {"output_type": "stream", "text": ["oi\n"]},
            ]},
            {"cell_type": "code", "source": "", "outputs": []},  # vazia → ignorada
        ],
        "nbformat": 4,
    }
    path = _write(".ipynb", json.dumps(nb).encode("utf-8"))
    try:
        text = extract.extract_text(path, "notebook.ipynb")
    finally:
        os.remove(path)
    assert "# Título" in text
    assert "texto md" in text
    assert "print('oi')" in text
    assert "[Saída]: oi" in text


def test_extract_image_placeholder_when_vision_disabled(monkeypatch):
    monkeypatch.setattr(settings, "vision_enabled", False)
    path = _write(".png", b"\x89PNG\r\n\x1a\n")  # bytes não precisam ser válidos: vision desligado
    try:
        text = extract.extract_text(path, "diagrama.png")
    finally:
        os.remove(path)
    assert text == "[Imagem: diagrama.png]"


def test_extract_pptx_text_by_slide(monkeypatch):
    monkeypatch.setattr(settings, "vision_enabled", False)  # sem chamadas de vision no teste
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # layout com título
    slide.shapes.title.text = "Arquitetura da Solução"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "GPON e integração de dados"

    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    try:
        text = extract.extract_text(path, "deck.pptx")
    finally:
        os.remove(path)
    assert "# Slide 1" in text
    assert "Arquitetura da Solução" in text
    assert "GPON e integração de dados" in text


def test_extract_ppt_legacy_unsupported():
    from app.domain.ingestion.errors import PermanentIngestionError

    path = _write(".ppt", b"legacy-ole-binary")
    try:
        raised = False
        try:
            extract.extract_text(path, "antigo.ppt")
        except PermanentIngestionError:
            raised = True
        assert raised
    finally:
        os.remove(path)
